from pptx import Presentation
from PIL import Image
import os
import io

def find_images_in_shape(shape):
    images = []
    if shape.shape_type == 13:
        images.append(shape)
    elif shape.shape_type == 6:
        for child in shape.shapes:
            images.extend(find_images_in_shape(child))
    return images

output_dir = "/Applications/HADEED CONSTRUCTION/public/images/projects"
os.makedirs(output_dir, exist_ok=True)

# === HEC.pptx ===
pptx1 = "/Users/rahulmedhe/Downloads/Hec.pptx"
prs1 = Presentation(pptx1)

# Page 6 - DEWA Solar Innovation Centre
slide6 = prs1.slides[5]
imgs6 = []
for shape in slide6.shapes:
    imgs6.extend(find_images_in_shape(shape))
print(f"Page 6 (DEWA): {len(imgs6)} images")
# Pick the largest image (main project photo)
imgs6_sorted = sorted(imgs6, key=lambda s: len(s.image.blob), reverse=True)
for i, img_shape in enumerate(imgs6_sorted):
    img_data = img_shape.image.blob
    img = Image.open(io.BytesIO(img_data))
    fname = f"dewa-solar-innovation-{i}.png"
    img.save(os.path.join(output_dir, fname), "PNG")
    print(f"  Saved {fname}: {img.size}")

# Page 7 - ICT 210
slide7 = prs1.slides[6]
imgs7 = []
for shape in slide7.shapes:
    imgs7.extend(find_images_in_shape(shape))
print(f"\nPage 7 (ICT 210): {len(imgs7)} images")
imgs7_sorted = sorted(imgs7, key=lambda s: len(s.image.blob), reverse=True)
for i, img_shape in enumerate(imgs7_sorted):
    img_data = img_shape.image.blob
    img = Image.open(io.BytesIO(img_data))
    fname = f"ict-logistics-park-{i}.png"
    img.save(os.path.join(output_dir, fname), "PNG")
    print(f"  Saved {fname}: {img.size}")

# === HEC Brief Presentation ===
pptx2 = "/Users/rahulmedhe/Downloads/HEC -  Brief Presentation.pptx"
prs2 = Presentation(pptx2)

brief_pages = {
    21: "dj53",
    24: "j199",
    41: "j128",
    42: "j129",
}

for page_num, code in brief_pages.items():
    slide = prs2.slides[page_num - 1]
    imgs = []
    for shape in slide.shapes:
        imgs.extend(find_images_in_shape(shape))
    
    # Get text for context
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    
    print(f"\nPage {page_num} ({code}): {len(imgs)} images")
    print(f"  Text snippets: {[t[:80] for t in texts if len(t) > 5]}")
    
    imgs_sorted = sorted(imgs, key=lambda s: len(s.image.blob), reverse=True)
    for i, img_shape in enumerate(imgs_sorted):
        img_data = img_shape.image.blob
        img = Image.open(io.BytesIO(img_data))
        fname = f"{code}-project-{i}.png"
        img.save(os.path.join(output_dir, fname), "PNG")
        print(f"  Saved {fname}: {img.size}")

print("\nDone! All images extracted.")
