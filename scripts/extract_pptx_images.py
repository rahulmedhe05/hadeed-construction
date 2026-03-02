"""
Extract all images from both HEC PowerPoint files.
Organizes images into categories based on slide content.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import hashlib

BASE_DIR = "/Applications/HADEED CONSTRUCTION/public/images"

def ensure_dir(path):
    os.makedirs(path, exist_ok=True)

def save_image(image_blob, ext, folder, name):
    ensure_dir(folder)
    filepath = os.path.join(folder, f"{name}.{ext}")
    with open(filepath, "wb") as f:
        f.write(image_blob)
    print(f"  Saved: {filepath} ({len(image_blob)} bytes)")
    return filepath

def extract_images_from_slide(slide, idx):
    """Extract all images from a slide, return list of (blob, ext, shape_name)"""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            blob = image.blob
            # Skip tiny placeholder images (< 5KB)
            if len(blob) < 5000:
                continue
            images.append((blob, ext, shape.name))
        # Check for images in group shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = child.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    blob = image.blob
                    if len(blob) < 5000:
                        continue
                    images.append((blob, ext, child.name))
    return images

def extract_brief_presentation():
    """Extract from HEC - Brief Presentation.pptx"""
    pptx_path = "/Users/rahulmedhe/Downloads/HEC -  Brief Presentation.pptx"
    prs = Presentation(pptx_path)
    
    print("=" * 60)
    print("EXTRACTING: HEC - Brief Presentation.pptx")
    print("=" * 60)
    
    # Slide 1: Company profile hero images
    print("\n--- Slide 1: Hero/Profile ---")
    images = extract_images_from_slide(prs.slides[0], 0)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/hero", f"hero-ppt-{i+1}")
    
    # Slide 2-3: Introduction images (project photos)
    print("\n--- Slides 2-3: Introduction ---")
    for slide_idx in [1, 2]:
        images = extract_images_from_slide(prs.slides[slide_idx], slide_idx)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/projects", f"intro-{slide_idx+1}-{i+1}")
    
    # Slide 5-6: Client Base images
    print("\n--- Slides 5-6: Client Base ---")
    for slide_idx in [4, 5]:
        images = extract_images_from_slide(prs.slides[slide_idx], slide_idx)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/clients", f"client-{slide_idx+1}-{i+1}")
    
    # Slide 7: Financial Resources (bank logos)
    print("\n--- Slide 7: Financial Resources ---")
    # Skip - these are small bank logos
    
    # Slide 8: Personnel table
    print("\n--- Slide 8: Personnel ---")
    images = extract_images_from_slide(prs.slides[7], 7)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/about", f"personnel-{i+1}")
    
    # Slide 9: Equipment
    print("\n--- Slide 9: Equipment ---")
    images = extract_images_from_slide(prs.slides[8], 8)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/equipment", f"equipment-{i+1}")
    
    # Slide 10: Awards
    print("\n--- Slide 10: Awards ---")
    images = extract_images_from_slide(prs.slides[9], 9)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/awards", f"award-{i+1}")
    
    # Slide 11: Quality Procedures
    print("\n--- Slide 11: Quality ---")
    images = extract_images_from_slide(prs.slides[10], 10)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/projects", f"quality-{i+1}")
    
    # Slide 12-13: Certificates
    print("\n--- Slides 12-13: Certificates ---")
    for slide_idx in [11, 12]:
        images = extract_images_from_slide(prs.slides[slide_idx], slide_idx)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/certificates", f"cert-{slide_idx+1}-{i+1}")
    
    # Slide 14-15: Licenses
    print("\n--- Slides 14-15: Licenses ---")
    for slide_idx in [13, 14]:
        images = extract_images_from_slide(prs.slides[slide_idx], slide_idx)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/licenses", f"license-{slide_idx+1}-{i+1}")
    
    # Slide 16: Client logos (HEC Projects header with many logos)
    print("\n--- Slide 16: Client Logos ---")
    # Extract ALL images from slide 16 including small ones (logos)
    for shape in prs.slides[15].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            blob = image.blob
            if len(blob) > 1000:  # Include small logos but skip tiny placeholders
                h = hashlib.md5(blob).hexdigest()[:8]
                save_image(blob, ext, f"{BASE_DIR}/client-logos", f"logo-{h}")
    
    # Slides 17-43: Major Projects
    print("\n--- Slides 17-43: Project Images ---")
    project_categories = {
        17: "major-projects",
        18: "ict-logistics-park",
        19: "business-hub",
        20: "renewable-energy",
        21: "al-taaqa-warehouse",
        22: "printing-publishing",
        23: "industrial-cable",
        24: "pharmaceutical",
        25: "aviation",
        26: "aviation",
        27: "completed-projects",
        28: "building-materials",
        29: "logistics-kimera",
        30: "automobile",
        31: "metal-industries",
        32: "oilfield",
        33: "paper-printing",
        34: "marine-ports",
        35: "military-police",
        36: "residential",
        37: "food-beverages",
        38: "water-stations",
        39: "education",
        40: "malls-showrooms",
        41: "energy-substation",
        42: "energy-substation",
        43: "qatar-project",
    }
    
    project_counter = 1
    for slide_idx in range(16, min(44, len(prs.slides))):
        slide = prs.slides[slide_idx]
        images = extract_images_from_slide(slide, slide_idx)
        cat = project_categories.get(slide_idx + 1, "other")
        
        # Get slide title text for context
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        
        if images:
            print(f"  Slide {slide_idx+1} ({cat}): {len(images)} images, text: {texts[:2]}")
            for i, (blob, ext, name) in enumerate(images):
                save_image(blob, ext, f"{BASE_DIR}/projects", f"ppt-project-{project_counter:02d}")
                project_counter += 1

def extract_experience_presentation():
    """Extract from HEC - EXPERIENCE_compressed.pptx"""
    pptx_path = "/Users/rahulmedhe/Downloads/HEC -  EXPERIENCE_compressed.pptx"
    prs = Presentation(pptx_path)
    
    print("\n" + "=" * 60)
    print("EXTRACTING: HEC - EXPERIENCE_compressed.pptx")
    print("=" * 60)
    
    # Slide 1: Cover
    print("\n--- Slide 1: Cover ---")
    images = extract_images_from_slide(prs.slides[0], 0)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/hero", f"exp-hero-{i+1}")
    
    # Slide 2: CEO - Founder image
    print("\n--- Slide 2: CEO/Founder ---")
    images = extract_images_from_slide(prs.slides[1], 1)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/about", f"founder-{i+1}")
    
    # Slide 3: Co-Founder
    print("\n--- Slide 3: Co-Founder ---")
    images = extract_images_from_slide(prs.slides[2], 2)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/about", f"cofounder-{i+1}")
    
    # Slide 4: Company overview
    print("\n--- Slide 4: Company Overview ---")
    images = extract_images_from_slide(prs.slides[3], 3)
    for i, (blob, ext, name) in enumerate(images):
        save_image(blob, ext, f"{BASE_DIR}/about", f"overview-{i+1}")
    
    # Slide 5: Another image
    print("\n--- Slide 5 ---")
    if len(prs.slides) > 4:
        images = extract_images_from_slide(prs.slides[4], 4)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/about", f"about-exp-{i+1}")
    
    # Slide 21: Org chart
    print("\n--- Slide 21: Org Chart ---")
    if len(prs.slides) > 20:
        images = extract_images_from_slide(prs.slides[20], 20)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/about", f"orgchart-{i+1}")
    
    # Slide 23: Markets & Services
    print("\n--- Slide 23: Markets ---")
    if len(prs.slides) > 22:
        images = extract_images_from_slide(prs.slides[22], 22)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/about", f"markets-{i+1}")
    
    # Slide 29: Printing factory
    print("\n--- Slide 29: Printing Factory ---")
    if len(prs.slides) > 28:
        images = extract_images_from_slide(prs.slides[28], 28)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/projects", f"exp-printing-{i+1}")
    
    # Slide 31: Awards
    print("\n--- Slide 31: Awards ---")
    if len(prs.slides) > 30:
        images = extract_images_from_slide(prs.slides[30], 30)
        for i, (blob, ext, name) in enumerate(images):
            save_image(blob, ext, f"{BASE_DIR}/awards", f"exp-award-{i+1}")
    
    # Slide 34-42: Project images
    print("\n--- Slides 34-42: Projects ---")
    exp_counter = 1
    for slide_idx in range(33, min(42, len(prs.slides))):
        images = extract_images_from_slide(prs.slides[slide_idx], slide_idx)
        if images:
            print(f"  Slide {slide_idx+1}: {len(images)} images")
            for i, (blob, ext, name) in enumerate(images):
                save_image(blob, ext, f"{BASE_DIR}/projects", f"exp-project-{exp_counter:02d}")
                exp_counter += 1

if __name__ == "__main__":
    extract_brief_presentation()
    extract_experience_presentation()
    
    # Print summary
    print("\n" + "=" * 60)
    print("EXTRACTION COMPLETE - Summary:")
    print("=" * 60)
    for folder in ["hero", "about", "clients", "client-logos", "certificates", "licenses", "awards", "equipment", "projects"]:
        path = os.path.join(BASE_DIR, folder)
        if os.path.exists(path):
            count = len([f for f in os.listdir(path) if os.path.isfile(os.path.join(path, f))])
            print(f"  {folder}: {count} images")
