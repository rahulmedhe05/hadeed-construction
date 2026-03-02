from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('/Users/rahulmedhe/Downloads/HEC -  Brief Presentation.pptx')

project_img_counter = 1
for slide_idx in range(16, min(44, len(prs.slides))):
    slide = prs.slides[slide_idx]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    
    img_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            if len(blob) >= 5000:
                img_count += 1

    if img_count > 0:
        print(f'Slide {slide_idx+1}: imgs={img_count}, ppt-project-{project_img_counter:02d} to ppt-project-{project_img_counter+img_count-1:02d}')
        print(f'  Title: {texts[:3]}')
        project_img_counter += img_count
    else:
        print(f'Slide {slide_idx+1}: NO IMAGES, texts={texts[:2]}')
