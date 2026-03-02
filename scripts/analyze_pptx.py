from pptx import Presentation
import os, json

results = {}

# Analyze Brief Presentation
pptx_path = "/Users/rahulmedhe/Downloads/HEC -  Brief Presentation.pptx"
prs = Presentation(pptx_path)

results["brief"] = []
for i, slide in enumerate(prs.slides):
    slide_data = {"slide": i+1, "texts": [], "images": []}
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_data["texts"].append(text)
        if shape.shape_type == 13:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            slide_data["images"].append({
                "content_type": image.content_type,
                "size": len(image.blob),
                "name": shape.name,
                "ext": ext
            })
    results["brief"].append(slide_data)

# Analyze Experience PPT
pptx_path2 = "/Users/rahulmedhe/Downloads/HEC -  EXPERIENCE_compressed.pptx"
prs2 = Presentation(pptx_path2)

results["experience"] = []
for i, slide in enumerate(prs2.slides):
    slide_data = {"slide": i+1, "texts": [], "images": []}
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_data["texts"].append(text)
        if shape.shape_type == 13:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            slide_data["images"].append({
                "content_type": image.content_type,
                "size": len(image.blob),
                "name": shape.name,
                "ext": ext
            })
    results["experience"].append(slide_data)

print(json.dumps(results, indent=2))
