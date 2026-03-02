from pptx import Presentation
from pptx.util import Inches, Emu
import os

# Check both PPTs for client logo slides
for pptx_file in ['HEC - Brief Presentation.pptx', 'HEC - EXPERIENCE_compressed.pptx']:
    if not os.path.exists(pptx_file):
        continue
    prs = Presentation(pptx_file)
    print(f'=== {pptx_file} ===')
    for slide_idx, slide in enumerate(prs.slides):
        texts = []
        img_count = 0
        img_sizes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == 13:  # Picture
                img_count += 1
                img = shape.image
                img_sizes.append(len(img.blob))
        text_preview = ' | '.join(texts[:5])
        if 'client' in text_preview.lower() or 'logo' in text_preview.lower() or img_count > 10:
            print(f'  Slide {slide_idx}: {img_count} images, text: {text_preview}')
            if img_sizes:
                print(f'    Image sizes: min={min(img_sizes)}, max={max(img_sizes)}, avg={sum(img_sizes)//len(img_sizes)}')
