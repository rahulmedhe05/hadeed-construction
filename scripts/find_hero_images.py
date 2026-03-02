"""
Identify the best large project images for hero slider use.
Check which slides they came from to know what they depict.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# Check the Brief PPT intro images (slides 1-3) and large project images
prs = Presentation('/Users/rahulmedhe/Downloads/HEC -  Brief Presentation.pptx')

print("=== BRIEF PPT - Slides with large images ===")
for slide_idx in range(len(prs.slides)):
    slide = prs.slides[slide_idx]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            if len(blob) > 200000:  # Only large images
                print(f"  Slide {slide_idx+1}: {shape.name}, {len(blob)} bytes, {shape.image.content_type}")
                print(f"    Context: {texts[:2]}")

print()

# Check the Experience PPT for large images
prs2 = Presentation('/Users/rahulmedhe/Downloads/HEC -  EXPERIENCE_compressed.pptx')

print("=== EXPERIENCE PPT - Slides with large images ===")
for slide_idx in range(len(prs2.slides)):
    slide = prs2.slides[slide_idx]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            if len(blob) > 200000:  # Only large images
                print(f"  Slide {slide_idx+1}: {shape.name}, {len(blob)} bytes, {shape.image.content_type}")
                print(f"    Context: {texts[:2]}")
